from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import os

def extract_ppt_contents(pptx_path, output_csv, image_output_dir="extracted_images"):
    prs = Presentation(pptx_path)
    records = []

    os.makedirs(image_output_dir, exist_ok=True)

    def process_shape(shape, slide_num, parent_group=None):
        # 1) Extract text
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            text = shape.text.strip()
            if text:
                records.append({
                    "slide_number": slide_num,
                    "content_type": "text",
                    "shape_name": shape.name,
                    "content": text,
                    "file_path": ""
                })

        # 2) Extract tables
        if hasattr(shape, "has_table") and shape.has_table:
            table = shape.table
            for row_idx, row in enumerate(table.rows, start=1):
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                records.append({
                    "slide_number": slide_num,
                    "content_type": "table",
                    "shape_name": shape.name,
                    "content": " | ".join(row_data),
                    "file_path": ""
                })

        # 3) Extract images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_ext = image.ext
            image_filename = f"slide_{slide_num}_{shape.shape_id}.{image_ext}"
            image_path = os.path.join(image_output_dir, image_filename)

            with open(image_path, "wb") as f:
                f.write(image.blob)

            records.append({
                "slide_number": slide_num,
                "content_type": "image",
                "shape_name": shape.name,
                "content": "",
                "file_path": image_path
            })

        # 4) Recurse into grouped shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in shape.shapes:
                process_shape(subshape, slide_num, parent_group=shape.name)

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            process_shape(shape, slide_num)

    df = pd.DataFrame(records)
    df.to_csv(output_csv, index=False, encoding="utf-8-sig")
    print(f"Extraction completed. CSV saved to: {output_csv}")
    print(f"Images saved in folder: {image_output_dir}")


if __name__ == "__main__":
    pptx_file = "C:/"           #replace with your respective path
    output_csv = "C:/"
    image_folder = "C:/"

    extract_ppt_contents(pptx_file, output_csv, image_folder)

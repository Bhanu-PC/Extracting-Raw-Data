#pip install python-pptx pandas (install pptx and pandas library)

from pptx import Presentation
import pandas as pd
import os
print(f"Pandas version: {pd.__version__}")
print(f"PPTX version: {Presentation.__module__}")
print("CSV file is saved in:", os.getcwd())

def extract_text_from_pptx(pptx_path, output_csv):
    prs = Presentation("C:/")  #replace your path
    rows = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        rows.append({
            "slide_number": slide_index,
            "text_content": " | ".join(slide_text)
        })

    df = pd.DataFrame(rows)
    df.to_csv(output_csv, index=False, encoding="utf-8-sig")
    print(f"CSV exported successfully: {output_csv}")


if __name__ == "__main__":
    pptx_file = "C:/"   # replace with your file path
    csv_file = "C:/Users/"
    extract_text_from_pptx(pptx_file, csv_file)
